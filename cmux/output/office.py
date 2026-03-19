"""PPTX and DOCX generation."""

from __future__ import annotations

import re
from pathlib import Path


def save_docx(content: str, output_dir: Path, filename_base: str) -> Path:
    """Convert markdown-ish content to a DOCX file."""
    from docx import Document

    doc = Document()
    lines = content.split("\n")

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif re.match(r"^\d+\.\s", stripped):
            text = re.sub(r"^\d+\.\s", "", stripped)
            doc.add_paragraph(text, style="List Number")
        else:
            doc.add_paragraph(stripped)

    path = output_dir / f"{filename_base}.docx"
    doc.save(str(path))
    return path


def save_pptx(content: str, output_dir: Path, filename_base: str) -> Path:
    """Convert structured content to a PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Parse content into slides
    slides_data = _parse_slides(content)

    for slide_data in slides_data:
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])

        if not bullets:
            # Title slide
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
        else:
            # Content slide
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet

    path = output_dir / f"{filename_base}.pptx"
    prs.save(str(path))
    return path


def _parse_slides(content: str) -> list[dict]:
    """Parse markdown-ish content into slide structures."""
    slides = []
    current_slide: dict | None = None

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("# ") or stripped.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            heading = re.sub(r"^#+\s*", "", stripped)
            current_slide = {"title": heading, "bullets": []}
        elif stripped.startswith("- ") or stripped.startswith("* "):
            if current_slide is None:
                current_slide = {"title": "", "bullets": []}
            current_slide["bullets"].append(stripped[2:])
        elif re.match(r"^\d+\.\s", stripped):
            if current_slide is None:
                current_slide = {"title": "", "bullets": []}
            text = re.sub(r"^\d+\.\s", "", stripped)
            current_slide["bullets"].append(text)

    if current_slide:
        slides.append(current_slide)

    if not slides:
        slides.append({"title": "Content", "bullets": [content[:200]]})

    return slides
